# -*- coding: utf-8 -*-
"""
文档生成技能
============

将 AI 生成的文本内容导出为 PDF、Word、PPT 等格式的文档。
支持 Markdown 转 PDF/Word，以及生成简单的 PPT。
"""

import os
import sys
import re
import time
import logging
from typing import Dict, Any

project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, project_root)

from skills import register_skill

logger = logging.getLogger(__name__)

OUTPUT_DIR = os.path.join(project_root, "data", "generated_docs")


def _ensure_output_dir():
    os.makedirs(OUTPUT_DIR, exist_ok=True)


def _sanitize_filename(filename: str) -> str:
    if not filename:
        return ""
    filename = re.sub(r'[\\/:*?"<>|]', "", filename)
    filename = re.sub(r'[^\w.\-]', "", filename, flags=re.UNICODE)
    return filename.strip()[:50]


def _extract_title_and_content(text: str):
    lines = text.strip().split("\n")
    title = ""
    content_start = 0
    for i, line in enumerate(lines):
        if line.strip().startswith("# "):
            title = line.strip().lstrip("# ").strip()
            content_start = i + 1
            break
    if not title:
        for i, line in enumerate(lines):
            stripped = line.strip()
            if stripped and not stripped.startswith("#") and not stripped.startswith("-") and not stripped.startswith("```"):
                title = stripped[:30]
                content_start = i
                break
    if not title:
        title = "AI生成的文档"
    content = "\n".join(lines[content_start:]).strip()
    return title, content


def _parse_markdown_to_structured(text: str):
    sections = []
    current_title = ""
    current_lines = []

    for line in text.split("\n"):
        if re.match(r'^#{1,3}\s+', line):
            if current_title or current_lines:
                sections.append({"title": current_title, "content": "\n".join(current_lines).strip()})
            current_title = re.sub(r'^#{1,3}\s+', '', line).strip()
            current_lines = []
        else:
            current_lines.append(line)

    if current_title or current_lines:
        sections.append({"title": current_title, "content": "\n".join(current_lines).strip()})

    if not sections:
        sections = [{"title": "", "content": text.strip()}]

    return sections


def _generate_pdf(title: str, content: str, filename: str) -> Dict[str, Any]:
    try:
        from fpdf import FPDF
    except ImportError:
        return {"success": False, "message": "PDF 生成需要安装 fpdf2：pip install fpdf2"}

    _ensure_output_dir()
    filepath = os.path.join(OUTPUT_DIR, filename)

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    font_path = os.path.join(project_root, "assets", "fonts", "SimHei.ttf")
    if os.path.exists(font_path):
        pdf.add_font("SimHei", "", font_path, uni=True)
        pdf.add_font("SimHei", "B", font_path, uni=True)
        font_name = "SimHei"
    else:
        font_name = "Helvetica"

    pdf.set_font(font_name, "B", 18)
    pdf.cell(0, 15, title, ln=True, align="C")
    pdf.ln(10)

    pdf.set_font(font_name, "", 11)
    for line in content.split("\n"):
        stripped = line.strip()
        if not stripped:
            pdf.ln(4)
            continue
        if stripped.startswith("# "):
            pdf.set_font(font_name, "B", 14)
            pdf.multi_cell(0, 8, stripped.lstrip("# "))
            pdf.ln(3)
            pdf.set_font(font_name, "", 11)
        elif stripped.startswith("## "):
            pdf.set_font(font_name, "B", 13)
            pdf.multi_cell(0, 7, stripped.lstrip("# "))
            pdf.ln(2)
            pdf.set_font(font_name, "", 11)
        elif stripped.startswith("### "):
            pdf.set_font(font_name, "B", 12)
            pdf.multi_cell(0, 7, stripped.lstrip("# "))
            pdf.ln(2)
            pdf.set_font(font_name, "", 11)
        elif stripped.startswith("- ") or stripped.startswith("* "):
            pdf.cell(5)
            pdf.multi_cell(0, 6, f"\u2022 {stripped[2:]}")
        elif re.match(r'^\d+\.\s+', stripped):
            pdf.multi_cell(0, 6, stripped)
        elif stripped.startswith("```"):
            continue
        else:
            pdf.multi_cell(0, 6, stripped)

    try:
        pdf.output(filepath)
    except Exception as e:
        logger.error(f"PDF 输出失败: {e}")
        return {"success": False, "message": f"PDF 生成失败: {e}"}

    file_size = os.path.getsize(filepath)
    return {
        "success": True,
        "filepath": filepath,
        "filename": filename,
        "file_size": file_size,
    }


def _generate_word(title: str, content: str, filename: str) -> Dict[str, Any]:
    try:
        from docx import Document
        from docx.shared import Pt, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        return {"success": False, "message": "Word 生成需要安装 python-docx：pip install python-docx"}

    _ensure_output_dir()
    filepath = os.path.join(OUTPUT_DIR, filename)

    doc = Document()

    title_para = doc.add_heading(title, level=0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    for line in content.split("\n"):
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith("# ") and not stripped.startswith("## "):
            continue
        if stripped.startswith("### "):
            doc.add_heading(stripped.lstrip("# "), level=3)
        elif stripped.startswith("## "):
            doc.add_heading(stripped.lstrip("# "), level=2)
        elif stripped.startswith("# "):
            doc.add_heading(stripped.lstrip("# "), level=1)
        elif stripped.startswith("- ") or stripped.startswith("* "):
            doc.add_paragraph(stripped[2:], style="List Bullet")
        elif re.match(r'^\d+\.\s+', stripped):
            doc.add_paragraph(re.sub(r'^\d+\.\s+', '', stripped), style="List Number")
        elif stripped.startswith("```"):
            continue
        else:
            para = doc.add_paragraph(stripped)
            for run in para.runs:
                run.font.size = Pt(11)

    try:
        doc.save(filepath)
    except Exception as e:
        logger.error(f"Word 输出失败: {e}")
        return {"success": False, "message": f"Word 生成失败: {e}"}

    file_size = os.path.getsize(filepath)
    return {
        "success": True,
        "filepath": filepath,
        "filename": filename,
        "file_size": file_size,
    }


def _generate_ppt(title: str, content: str, filename: str) -> Dict[str, Any]:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return {"success": False, "message": "PPT 生成需要安装 python-pptx：pip install python-pptx"}

    _ensure_output_dir()
    filepath = os.path.join(OUTPUT_DIR, filename)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]

    def _add_title_slide(prs, title_text):
        slide = prs.slides.add_slide(slide_layout)
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(11.333)
        height = Inches(2)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.alignment = PP_ALIGN.CENTER
        return slide

    def _add_content_slide(prs, slide_title, slide_content):
        slide = prs.slides.add_slide(slide_layout)

        left = Inches(0.8)
        top = Inches(0.5)
        width = Inches(11.733)
        height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

        left = Inches(0.8)
        top = Inches(1.8)
        width = Inches(11.733)
        height = Inches(5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        bullet_items = []
        for line in slide_content.split("\n"):
            stripped = line.strip()
            if not stripped:
                continue
            clean = stripped
            for prefix in ["- ", "* ", "  - ", "  * "]:
                if clean.startswith(prefix):
                    clean = clean[len(prefix):]
                    break
            if re.match(r'^\d+\.\s+', clean):
                clean = re.sub(r'^\d+\.\s+', '', clean)
            bullet_items.append(clean)

        for i, item in enumerate(bullet_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
            p.space_after = Pt(8)
            p.level = 0

        return slide

    sections = _parse_markdown_to_structured(content)

    _add_title_slide(prs, title)

    for section in sections:
        if section["title"]:
            _add_content_slide(prs, section["title"], section["content"])
        elif section["content"]:
            content_lines = [l for l in section["content"].split("\n") if l.strip()]
            if content_lines:
                slide_title = content_lines[0][:40]
                slide_content = "\n".join(content_lines[1:]) if len(content_lines) > 1 else content_lines[0]
                _add_content_slide(prs, slide_title, slide_content)

    try:
        prs.save(filepath)
    except Exception as e:
        logger.error(f"PPT 输出失败: {e}")
        return {"success": False, "message": f"PPT 生成失败: {e}"}

    file_size = os.path.getsize(filepath)
    return {
        "success": True,
        "filepath": filepath,
        "filename": filename,
        "file_size": file_size,
    }


_FORMAT_GENERATORS = {
    "pdf": _generate_pdf,
    "word": _generate_word,
    "ppt": _generate_ppt,
}

_FORMAT_EXTENSIONS = {
    "pdf": ".pdf",
    "word": ".docx",
    "ppt": ".pptx",
}

_FORMAT_NAMES = {
    "pdf": "PDF",
    "word": "Word",
    "ppt": "PPT",
}

_FORMAT_KEYWORDS = {
    "pdf": ["pdf", "PDF"],
    "word": ["word", "doc", "docx", "Word", "文档"],
    "ppt": ["ppt", "pptx", "PPT", "幻灯片", "演示文稿", "汇报"],
}


def _detect_format(text: str) -> str:
    lower = text.lower()
    for fmt, keywords in _FORMAT_KEYWORDS.items():
        for kw in keywords:
            if kw in lower:
                return fmt
    return "pdf"


@register_skill(
    skill_id="doc_generator",
    name="文档生成",
    description="将 AI 内容导出为 PDF、Word 或 PPT 文档并下载",
    triggers=["生成文档", "导出文档", "生成PDF", "导出PDF", "生成Word", "导出Word",
              "生成PPT", "导出PPT", "生成pdf", "生成word", "生成ppt",
              "转为文档", "转为PDF", "转为Word", "转为PPT",
              "制作文档", "制作PPT", "制作Word",
              "下载文档", "下载PDF", "下载报告",
              "生成报告", "导出报告"],
    icon="file",
    examples=[
        "帮我把以下内容生成PDF文档：...",
        "生成一份PPT关于AI发展趋势",
        "导出为Word文档",
    ],
    tool_schema={
        "type": "function",
        "function": {
            "name": "generate_document",
            "description": "将文本内容生成指定格式的文档（PDF/Word/PPT）。当用户要求生成、导出、下载文档时使用此工具。",
            "parameters": {
                "type": "object",
                "properties": {
                    "prompt": {
                        "type": "string",
                        "description": "要生成文档的主题或完整内容。如果用户提供了完整内容则直接使用；如果只给了主题，则先生成内容再导出。",
                    },
                    "format": {
                        "type": "string",
                        "description": "文档格式",
                        "enum": ["pdf", "word", "ppt"],
                    },
                },
                "required": ["prompt"],
            },
        },
    },
)
def handle_doc_generator(user_input: str, context: Dict[str, Any]) -> Dict[str, Any]:
    text = user_input
    for trigger in ["生成文档", "导出文档", "转为文档", "制作文档", "下载文档",
                     "生成PDF", "导出PDF", "转为PDF", "下载PDF",
                     "生成Word", "导出Word", "转为Word", "制作Word",
                     "生成PPT", "导出PPT", "转为PPT", "制作PPT", "下载报告",
                     "生成报告", "导出报告",
                     "生成pdf", "生成word", "生成ppt",
                     "帮我", "请帮我", "把", "一下", "的", "文档"]:
        text = text.replace(trigger, "").strip()

    tool_args = context.get("tool_args", {}) if context else {}
    fmt = tool_args.get("format") or _detect_format(user_input)

    prompt_text = tool_args.get("prompt") or text
    if not prompt_text:
        prompt_text = user_input

    has_structured_content = any(c in prompt_text for c in ["#", "\n\n", "- ", "1.", "##"])

    if not has_structured_content:
        try:
            from core.agents import get_agent_manager
            from core.model_router import build_llm_for_agent
            manager = get_agent_manager(os.path.join(project_root, "data"))
            agent_config = manager.get_default_agent()
            llm = build_llm_for_agent(agent_config)

            if fmt == "ppt":
                content_prompt = (
                    f"请就以下主题生成一份 PPT 大纲内容，用 Markdown 格式输出。\n"
                    f"每个一级标题（# ）作为幻灯片页的标题，其下方的内容作为该页的要点（用 - 列表）。\n"
                    f"生成 5-8 页幻灯片的内容，每页 3-5 个要点。\n\n"
                    f"主题：{prompt_text}"
                )
            else:
                content_prompt = (
                    f"请就以下主题生成一份结构化的文档内容，用 Markdown 格式输出。\n"
                    f"包含标题、章节、列表等格式。\n\n"
                    f"主题：{prompt_text}"
                )

            resp = llm.call(messages=[
                {"role": "system", "content": "你是一个专业的文档撰写助手，擅长生成结构清晰的内容。"},
                {"role": "user", "content": content_prompt},
            ])
            prompt_text = str(resp)
        except Exception as e:
            logger.warning(f"AI 内容生成失败，使用原始输入: {e}")

    title, content = _extract_title_and_content(prompt_text)
    if not content:
        content = prompt_text

    safe_name = _sanitize_filename(title) or "document"
    ext = _FORMAT_EXTENSIONS.get(fmt, ".pdf")
    filename = f"{safe_name}_{int(time.time())}{ext}"

    generator = _FORMAT_GENERATORS.get(fmt)
    if not generator:
        return {"success": False, "message": f"不支持的文档格式: {fmt}"}

    result = generator(title, content, filename)

    if result["success"]:
        size_kb = result["file_size"] / 1024
        download_url = f"/api/files/docs/{filename}"
        format_name = _FORMAT_NAMES.get(fmt, fmt.upper())
        msg = (
            f"**{format_name} 文档已生成！**\n\n"
            f"**标题：** {title}\n\n"
            f"**大小：** {size_kb:.1f} KB\n\n"
            f"[点击下载 {format_name} 文档]({download_url})"
        )
        return {"success": True, "message": msg, "download_url": download_url, "filename": filename}
    else:
        return {"success": False, "message": f"文档生成失败: {result['message']}"}
